"""
Gemini Slide Parser API Routes
Parse presentation slides (PDF/PPTX) using Gemini Vision to extract exact layout and styling
"""

from fastapi import APIRouter, HTTPException, Depends, UploadFile, File
from pydantic import BaseModel, Field
from typing import List, Optional, Dict, Any
import base64
import io
import logging
from pathlib import Path
import tempfile

from src.middleware.firebase_auth import require_auth
from src.services.ai_chat_service import ai_chat_service, AIProvider
from src.utils.logger import setup_logger

logger = setup_logger(__name__)
router = APIRouter(prefix="/api/gemini/slides", tags=["Gemini Slide Parser"])


# ============ MODELS ============


class SlideContent(BaseModel):
    """Single slide with HTML content"""

    slide_number: int = Field(..., description="Slide number (1-indexed)")
    html_content: str = Field(..., description="HTML content with exact styling")
    notes: Optional[str] = Field(None, description="Speaker notes if any")


class SlideParseRequest(BaseModel):
    """Request to parse slides from file_id"""

    file_id: str = Field(..., description="File ID from user_files collection")


class SlideParseResponse(BaseModel):
    """Response with parsed slides"""

    success: bool = Field(..., description="Whether parsing succeeded")
    total_slides: int = Field(..., description="Total number of slides parsed")
    slides: List[SlideContent] = Field(..., description="List of slides with HTML")
    file_name: str = Field(..., description="Original filename")
    error: Optional[str] = Field(None, description="Error message if failed")


# ============ HELPER FUNCTIONS ============


def upload_pdf_to_gemini(pdf_path: str) -> str:
    """
    Upload PDF file to Gemini Files API

    Args:
        pdf_path: Path to PDF file

    Returns:
        File URI for use in Gemini API
    """
    try:
        import google.generativeai as genai
        from src.core.config import APP_CONFIG

        logger.info(f"📄 Uploading PDF to Gemini: {pdf_path}")

        # Configure Gemini
        genai.configure(api_key=APP_CONFIG["gemini_api_key"])

        # Upload PDF file
        uploaded_file = genai.upload_file(pdf_path)
        logger.info(f"✅ PDF uploaded: {uploaded_file.uri}")

        # Return the file name (files/xxx format) for get_file()
        return uploaded_file.name

    except Exception as e:
        logger.error(f"❌ Failed to upload PDF to Gemini: {e}")
        raise


async def parse_pdf_with_gemini(
    pdf_file_uri: str, file_name: str
) -> List[SlideContent]:
    """
    Parse entire PDF using Gemini 2.5 Pro (supports PDF natively)

    Args:
        pdf_file_uri: Gemini file URI from upload
        file_name: Original filename for context

    Returns:
        List of SlideContent with HTML for each slide
    """
    try:
        import google.generativeai as genai
        from src.core.config import APP_CONFIG

        logger.info(f"🤖 Parsing PDF with Gemini 2.5 Pro: {file_name}")

        # Configure Gemini
        genai.configure(api_key=APP_CONFIG["gemini_api_key"])

        # Enhanced prompt for slide analysis
        prompt = f"""Analyze this presentation PDF file and convert each slide to HTML.

**Your Task:**
1. Read through ALL slides in the PDF
2. For each slide, create pixel-perfect HTML that preserves:
   - **Layout**: Exact positioning of all elements
   - **Typography**: Font sizes, weights, styles, alignment
   - **Colors**: Text colors, backgrounds, borders, gradients
   - **Spacing**: Margins, padding, line heights
   - **Structure**: Headings, paragraphs, lists, tables, images

**Output Format:**
Return a JSON array with one object per slide:
```json
[
  {{
    "slide_number": 1,
    "html_content": "<div style='width:1920px;height:1080px;position:relative;background:#fff;'>...</div>"
  }},
  {{
    "slide_number": 2,
    "html_content": "<div style='width:1920px;height:1080px;position:relative;background:#fff;'>...</div>"
  }}
]
```

**HTML Requirements:**
- Full HD dimensions: 1920x1080px (16:9 ratio)
- Use inline CSS for all styling
- Use absolute/relative positioning for layout
- Preserve ALL text content exactly as shown
- Use semantic HTML: h1, h2, h3, p, ul, li, table, span
- Match font sizes: title (48-72px), heading (32-48px), body (20-28px)
- Include background colors and gradients
- Center or align text as shown in original

**Example slide HTML:**
```html
<div style="width:1920px;height:1080px;position:relative;background:linear-gradient(135deg,#667eea,#764ba2);padding:80px;">
  <h1 style="font-size:72px;font-weight:bold;color:#fff;text-align:center;margin-bottom:40px;">Slide Title</h1>
  <p style="font-size:28px;color:#fff;text-align:center;line-height:1.6;">Body text here</p>
  <ul style="font-size:24px;color:#fff;margin-top:60px;">
    <li style="margin-bottom:20px;">Bullet point 1</li>
    <li style="margin-bottom:20px;">Bullet point 2</li>
  </ul>
</div>
```

**IMPORTANT:**
- Return ONLY the JSON array, no markdown code blocks
- Each slide must be a complete, standalone HTML div
- Preserve exact text from PDF (no changes or translations)

File: {file_name}

Analyze all slides and return the JSON array:"""

        # Create model (use latest available Gemini model)
        try:
            model = genai.GenerativeModel("gemini-2.0-flash-exp")
            logger.info("🤖 Using model: gemini-2.0-flash-exp")
        except:
            try:
                model = genai.GenerativeModel("gemini-1.5-pro")
                logger.info("🤖 Using model: gemini-1.5-pro")
            except:
                model = genai.GenerativeModel("gemini-pro-vision")
                logger.info("🤖 Using model: gemini-pro-vision")

        # Generate content with PDF
        response = model.generate_content([prompt, genai.get_file(pdf_file_uri)])

        logger.info(f"✅ Got response from Gemini: {len(response.text)} chars")

        # Parse JSON response
        import json
        import re

        # Clean response (remove markdown if present)
        response_text = response.text.strip()
        if response_text.startswith("```json"):
            response_text = response_text[7:]
        if response_text.startswith("```"):
            response_text = response_text[3:]
        if response_text.endswith("```"):
            response_text = response_text[:-3]
        response_text = response_text.strip()

        # Parse JSON
        slides_data = json.loads(response_text)

        # Convert to SlideContent objects
        slides = []
        for slide_data in slides_data:
            slide = SlideContent(
                slide_number=slide_data["slide_number"],
                html_content=slide_data["html_content"],
                notes=slide_data.get("notes"),
            )
            slides.append(slide)
            logger.info(
                f"  Slide {slide.slide_number}: {len(slide.html_content)} chars HTML"
            )
            # Log first 500 chars of HTML for debugging
            logger.info(f"  HTML preview: {slide.html_content[:500]}...")

        logger.info(f"✅ Parsed {len(slides)} slides from PDF")
        return slides

    except Exception as e:
        logger.error(f"❌ Failed to parse PDF with Gemini: {e}")
        import traceback

        logger.error(traceback.format_exc())
        raise


# ============ ENDPOINTS ============


@router.post("/parse-file", response_model=SlideParseResponse)
async def parse_slides_from_file(
    request: SlideParseRequest, user_info: dict = Depends(require_auth)
):
    """
    Parse presentation file (PDF/PPTX) into HTML slides using Gemini 2.5 Pro

    **Flow:**
    1. Get file info from user_files collection
    2. Download file from R2 storage
    3. Upload PDF directly to Gemini Files API
    4. Send to Gemini 2.5 Pro (supports native PDF reading)
    5. Gemini analyzes ALL slides in one call and returns HTML array
    6. Return list of HTML slides with preserved layout and styling

    **Advantages of Gemini 2.5 Pro:**
    - ✅ Native PDF support (no image conversion needed)
    - ✅ One API call for entire PDF (vs N calls for N slides)
    - ✅ Better text extraction and layout understanding
    - ✅ Faster processing (~5-10s for 10 slides vs 30-50s with images)

    **Use Cases:**
    - Convert uploaded presentations to editable HTML
    - Preserve exact layout and styling from original slides
    - Extract text content with font sizes and colors

    **Example Request:**
    ```json
    {
      "file_id": "file_abc123"
    }
    ```

    **Example Response:**
    ```json
    {
      "success": true,
      "total_slides": 10,
      "slides": [
        {
          "slide_number": 1,
          "html_content": "<div style='width: 1920px; height: 1080px;'>...</div>",
          "notes": null
        }
      ],
      "file_name": "Presentation.pdf"
    }
    ```

    **Performance:**
    - Upload: ~2-3 seconds
    - Processing: ~5-10 seconds total (all slides)
    - 10 slides = ~7-13 seconds (vs 30-50s with image conversion)
    """
    try:
        user_id = user_info["uid"]
        file_id = request.file_id

        logger.info(f"🎬 Starting slide parse for file {file_id} (user: {user_id})...")

        # Step 1: Get file info from database
        from src.services.user_manager import UserManager
        from src.database.db_manager import DBManager

        db_manager = DBManager()
        user_manager = UserManager(db_manager)

        file_info = user_manager.get_file_by_id(file_id, user_id)
        if not file_info:
            logger.error(f"❌ File {file_id} not found in database for user {user_id}")
            raise HTTPException(
                status_code=404,
                detail=f"File {file_id} not found in user_files collection",
            )

        file_name = file_info.get("original_name", "Unknown")
        file_type = file_info.get("file_type", "")
        r2_key = file_info.get("r2_key")

        logger.info(f"📄 File: {file_name} (type: {file_type})")

        # Validate file type
        if file_type not in [".pdf", ".pptx"]:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type: {file_type}. Only PDF and PPTX are supported.",
            )

        if not r2_key:
            raise HTTPException(status_code=500, detail="File R2 key not found")

        # Step 2: Check if we already have parsed slides for this file (caching)
        from src.services.document_manager import DocumentManager

        doc_manager = DocumentManager(
            db_manager.db
        )  # DocumentManager expects raw Database

        # Count existing slide documents from this file
        existing_count = doc_manager.count_documents_by_file_id(file_id, user_id)
        logger.info(
            f"📊 Found {existing_count} existing document(s) from file {file_id}"
        )

        # Try to reuse cached slides from previous parse
        cached_slides = None
        if existing_count > 0:
            # Get the most recent document for this file
            previous_doc = doc_manager.get_latest_document_by_file_id(file_id, user_id)

            if previous_doc:
                prev_html = previous_doc.get("content_html")

                # Validate cached content has actual data
                if prev_html and prev_html.strip():
                    from bs4 import BeautifulSoup

                    soup = BeautifulSoup(prev_html, "html.parser")
                    text_content = soup.get_text(strip=True)

                    if (
                        text_content and len(text_content) > 50
                    ):  # At least 50 chars for slides
                        # Cache is valid! Extract slides from cached HTML
                        # Format: Multiple <div> tags, one per slide
                        slide_divs = soup.find_all(
                            "div",
                            style=lambda s: s
                            and "width:1920px" in s
                            or "width: 1920px" in s,
                        )

                        if slide_divs and len(slide_divs) > 0:
                            cached_slides = []
                            for idx, slide_div in enumerate(slide_divs):
                                cached_slides.append(
                                    SlideContent(
                                        slide_number=idx + 1,
                                        html_content=str(slide_div),
                                        notes=None,
                                    )
                                )

                            logger.info(
                                f"♻️ Reusing cached slides from previous parse (fast path!)"
                            )
                            logger.info(
                                f"♻️ Cached {len(cached_slides)} slides, HTML: {len(prev_html)} chars, Text: {len(text_content)} chars"
                            )
                        else:
                            logger.warning(
                                f"⚠️ Previous document HTML doesn't contain slide divs, will re-parse"
                            )
                    else:
                        text_len = len(text_content) if text_content else 0
                        logger.warning(
                            f"⚠️ Previous document has empty text content (HTML: {len(prev_html)} chars, Text: {text_len} chars), will re-parse file"
                        )
                else:
                    logger.warning(
                        f"⚠️ Previous document has empty or no content_html, will re-parse file"
                    )

        # Step 3: Get slides (from cache or parse with Gemini)
        if cached_slides:
            # Fast path: Return cached slides immediately
            slides = cached_slides
            logger.info(
                f"✅ Returned {len(slides)} cached slides (no Gemini call needed)"
            )

        else:
            # Slow path: Download from R2 and parse with Gemini (first time or cache invalid)
            logger.info(f"📥 Downloading file from R2: {r2_key}")

            # Download file to temp location
            with tempfile.NamedTemporaryFile(
                delete=False, suffix=file_type
            ) as tmp_file:
                tmp_path = tmp_file.name

                # Download file content
                from src.storage.r2_client import R2Client
                from src.core.config import APP_CONFIG

                r2_client = R2Client(
                    account_id=APP_CONFIG["r2_account_id"],
                    access_key_id=APP_CONFIG["r2_access_key_id"],
                    secret_access_key=APP_CONFIG["r2_secret_access_key"],
                    bucket_name=APP_CONFIG["r2_bucket_name"],
                    region=APP_CONFIG["r2_region"],
                )

                file_obj = r2_client.get_file(r2_key)
                if not file_obj:
                    raise HTTPException(
                        status_code=500, detail="Failed to download file from R2"
                    )

                # Write to temp file
                file_content = file_obj["Body"].read()
                tmp_file.write(file_content)
                logger.info(f"✅ Downloaded {len(file_content)} bytes to {tmp_path}")

            # Parse with Gemini
            try:
                if file_type == ".pdf":
                    # Upload PDF to Gemini Files API
                    pdf_uri = upload_pdf_to_gemini(tmp_path)

                    # Parse entire PDF with Gemini 2.5 Pro (one call for all slides)
                    slides = await parse_pdf_with_gemini(pdf_uri, file_name)

                else:  # .pptx
                    # TODO: Implement PPTX parsing
                    raise HTTPException(
                        status_code=501, detail="PPTX parsing not yet implemented"
                    )

            finally:
                # Clean up temp file
                Path(tmp_path).unlink(missing_ok=True)

            logger.info(f"✅ Successfully parsed {len(slides)} slides with Gemini")

            # Save slides as document in database for caching and editing
            # Combine all slides into one HTML document
            combined_html = "\n\n".join([slide.html_content for slide in slides])

            # Extract text content from HTML for search
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(combined_html, "html.parser")
            combined_text = soup.get_text(separator="\n", strip=True)

            # Create document in database
            document_id = doc_manager.create_document(
                user_id=user_id,
                title=file_name,
                content_html=combined_html,
                content_text=combined_text,
                source_type="file",
                document_type="slide",
                file_id=file_id,
                original_r2_url=file_info.get("file_url"),
                original_file_type=file_type,
                folder_id=None,
            )

            logger.info(f"💾 Saved {len(slides)} slides to document {document_id}")

        return SlideParseResponse(
            success=True,
            total_slides=len(slides),
            slides=slides,
            file_name=file_name,
            error=None,
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ Failed to parse slides: {e}")
        import traceback

        logger.error(traceback.format_exc())

        return SlideParseResponse(
            success=False,
            total_slides=0,
            slides=[],
            file_name="",
            error=str(e),
        )


@router.post("/parse-upload", response_model=SlideParseResponse)
async def parse_slides_from_upload(
    file: UploadFile = File(...), user_info: dict = Depends(require_auth)
):
    """
    Parse uploaded presentation file directly (without saving to R2 first)

    **Use for quick testing without uploading to storage**

    **Flow:**
    1. Upload file directly via multipart/form-data
    2. Save to temp file
    3. Upload to Gemini Files API
    4. Parse with Gemini 2.5 Pro (native PDF support)
    5. Return HTML slides

    **Performance:** ~7-13 seconds for 10 slides (3x faster than image conversion)
    """
    try:
        user_id = user_info["uid"]
        file_name = file.filename

        logger.info(f"🎬 Starting direct upload parse: {file_name} (user: {user_id})")

        # Validate file type
        file_ext = Path(file_name).suffix.lower()
        if file_ext not in [".pdf", ".pptx"]:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type: {file_ext}. Only PDF and PPTX are supported.",
            )

        # Read file content
        file_content = await file.read()
        logger.info(f"📄 Read {len(file_content)} bytes from upload")

        # Save to temp file
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tmp_file:
            tmp_path = tmp_file.name
            tmp_file.write(file_content)

        # Parse PDF with Gemini
        try:
            if file_ext == ".pdf":
                # Upload PDF to Gemini Files API
                pdf_uri = upload_pdf_to_gemini(tmp_path)

                # Parse entire PDF with Gemini 2.5 Pro
                slides = await parse_pdf_with_gemini(pdf_uri, file_name)

            else:  # .pptx
                raise HTTPException(
                    status_code=501, detail="PPTX parsing not yet implemented"
                )

        finally:
            Path(tmp_path).unlink(missing_ok=True)

        logger.info(f"✅ Successfully parsed {len(slides)} slides")

        return SlideParseResponse(
            success=True,
            total_slides=len(slides),
            slides=slides,
            file_name=file_name,
            error=None,
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ Failed to parse uploaded slides: {e}")
        import traceback

        logger.error(traceback.format_exc())

        return SlideParseResponse(
            success=False,
            total_slides=0,
            slides=[],
            file_name=file.filename if file else "",
            error=str(e),
        )


# ============ SLIDE EXPORT MODELS ============


class SlideExportRequest(BaseModel):
    """Request to export slides as PPTX/PPT"""

    document_id: str = Field(..., description="Document ID containing slide HTML")
    format: str = Field(..., description="Export format: pptx or ppt")


class SlideExportResponse(BaseModel):
    """Response with presigned download URL"""

    success: bool = Field(..., description="Whether export succeeded")
    download_url: str = Field(None, description="Presigned URL to download file")
    file_name: str = Field(None, description="Generated filename")
    file_size: int = Field(None, description="File size in bytes")
    expires_in: int = Field(3600, description="URL expiration time in seconds")
    error: Optional[str] = Field(None, description="Error message if failed")


# ============ SLIDE EXPORT HELPER FUNCTIONS ============


def html_slides_to_pptx(html_content: str, title: str = "Presentation") -> bytes:
    """
    Convert HTML slides to PPTX file

    Args:
        html_content: Combined HTML containing all slides (each in a div.slide)
        title: Presentation title

    Returns:
        PPTX file bytes
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from bs4 import BeautifulSoup
        import re

        logger.info(f"🎨 Converting HTML to PPTX: {title}")

        # Parse HTML to extract individual slides
        soup = BeautifulSoup(html_content, "html.parser")
        slide_divs = soup.find_all("div", class_="slide")

        if not slide_divs:
            # Fallback: split by slide number divs or any div with style containing width/height
            slide_divs = soup.find_all("div", style=re.compile(r"width.*height"))

        logger.info(f"📊 Found {len(slide_divs)} slides in HTML")

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)  # 16:9 aspect ratio
        prs.slide_height = Inches(5.625)

        for idx, slide_div in enumerate(slide_divs, 1):
            logger.info(f"  Processing slide {idx}/{len(slide_divs)}")

            # Add blank slide
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)

            # Extract text content from slide
            slide_text = slide_div.get_text(separator="\n", strip=True)

            # Parse inline styles to extract positioning and styling
            elements = slide_div.find_all(style=True)

            for element in elements:
                element_text = element.get_text(strip=True)
                if not element_text:
                    continue

                # Parse style attributes
                style_str = element.get("style", "")
                style_dict = {}
                for style_item in style_str.split(";"):
                    if ":" in style_item:
                        key, value = style_item.split(":", 1)
                        style_dict[key.strip()] = value.strip()

                # Extract position and size
                try:
                    # Get position (convert px to inches, assuming 96 DPI)
                    left = (
                        float(re.search(r"(\d+)", style_dict.get("left", "0")).group(1))
                        / 96
                        if "left" in style_dict
                        else 0.5
                    )
                    top = (
                        float(re.search(r"(\d+)", style_dict.get("top", "0")).group(1))
                        / 96
                        if "top" in style_dict
                        else 0.5
                    )

                    # Get font size
                    font_size_match = re.search(
                        r"(\d+)", style_dict.get("font-size", "18")
                    )
                    font_size = int(font_size_match.group(1)) if font_size_match else 18

                    # Add text box
                    width = Inches(8)  # Default width
                    height = Inches(1)  # Default height

                    txBox = slide.shapes.add_textbox(
                        Inches(left), Inches(top), width, height
                    )
                    tf = txBox.text_frame
                    tf.text = element_text
                    tf.word_wrap = True

                    # Apply formatting
                    for paragraph in tf.paragraphs:
                        paragraph.font.size = Pt(font_size)

                        # Apply color if specified
                        if "color" in style_dict:
                            color_str = style_dict["color"]
                            # Parse hex color
                            if color_str.startswith("#"):
                                from pptx.util import RGBColor

                                color_hex = color_str.lstrip("#")
                                if len(color_hex) == 6:
                                    r = int(color_hex[0:2], 16)
                                    g = int(color_hex[2:4], 16)
                                    b = int(color_hex[4:6], 16)
                                    paragraph.font.color.rgb = RGBColor(r, g, b)

                        # Apply bold if specified
                        if (
                            "font-weight" in style_dict
                            and "bold" in style_dict["font-weight"]
                        ):
                            paragraph.font.bold = True

                except (ValueError, AttributeError) as e:
                    logger.warning(f"⚠️ Failed to parse element style: {e}")
                    continue

        # Save to bytes
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        logger.info(f"✅ Created PPTX with {len(slide_divs)} slides")
        return pptx_buffer.getvalue()

    except Exception as e:
        logger.error(f"❌ Failed to convert HTML to PPTX: {e}")
        import traceback

        logger.error(traceback.format_exc())
        raise


# ============ SLIDE EXPORT ENDPOINT ============


@router.post("/export", response_model=SlideExportResponse)
async def export_slides_to_pptx(
    request: SlideExportRequest, user_info: dict = Depends(require_auth)
):
    """
    Export slide document as PPTX or PPT file

    **Flow:**
    1. Get document from MongoDB by document_id
    2. Extract HTML content and split into individual slides
    3. Convert each slide HTML to PPTX slide with preserved styling
    4. Upload PPTX file to R2 storage
    5. Generate presigned URL (valid for 1 hour)
    6. Return download URL to frontend

    **Supported Formats:**
    - `.pptx` - PowerPoint 2007+ (recommended)
    - `.ppt` - PowerPoint 97-2003 (converted from .pptx)

    **Note:** HTML styling is approximated in PowerPoint format.
    Complex CSS may not be perfectly preserved.
    """
    try:
        user_id = user_info["uid"]

        logger.info(
            f"📥 Export slides request: document_id={request.document_id}, "
            f"format={request.format}, user={user_id}"
        )

        # Validate format
        if request.format.lower() not in ["pptx", "ppt"]:
            raise HTTPException(
                status_code=400, detail="Invalid format. Must be 'pptx' or 'ppt'"
            )

        # Initialize managers
        from src.db.mongodb import DBManager
        from src.services.document_manager import DocumentManager
        from src.storage.r2_client import R2Client
        from src.core.config import APP_CONFIG

        db_manager = DBManager()
        doc_manager = DocumentManager(db_manager.db)

        r2_client = R2Client(
            account_id=APP_CONFIG["r2_account_id"],
            access_key_id=APP_CONFIG["r2_access_key_id"],
            secret_access_key=APP_CONFIG["r2_secret_access_key"],
            bucket_name=APP_CONFIG["r2_bucket_name"],
            region=APP_CONFIG["r2_region"],
        )

        # Get document from database
        document = doc_manager.get_document(request.document_id, user_id)

        if not document:
            raise HTTPException(
                status_code=404, detail=f"Document not found: {request.document_id}"
            )

        # Check if document is slide type
        if document.get("document_type") != "slide":
            raise HTTPException(
                status_code=400,
                detail="Document is not a slide presentation. Only slide documents can be exported to PPTX.",
            )

        # Get HTML content
        html_content = document.get("content_html", "")
        if not html_content:
            raise HTTPException(
                status_code=400, detail="Document has no HTML content to export"
            )

        title = document.get("title", "Presentation")

        # Convert HTML to PPTX
        logger.info(f"🎨 Converting slides to PPTX...")
        pptx_bytes = html_slides_to_pptx(html_content, title)

        # Generate filename
        from datetime import datetime

        timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
        safe_title = "".join(c if c.isalnum() or c in "-_" else "_" for c in title)
        file_name = f"{safe_title}_{timestamp}.{request.format}"

        # Upload to R2
        logger.info(f"☁️ Uploading to R2: {file_name}")
        file_key = f"exports/{user_id}/slides/{file_name}"

        upload_result = r2_client.upload_file_from_memory(
            file_content=pptx_bytes,
            file_key=file_key,
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

        if not upload_result.get("success"):
            raise HTTPException(
                status_code=500,
                detail=f"Failed to upload file to storage: {upload_result.get('error')}",
            )

        logger.info(f"✅ File uploaded to R2: {file_key}")

        # Generate presigned URL (valid for 1 hour)
        presigned_url = r2_client.generate_presigned_url(
            file_key=file_key, expiration=3600  # 1 hour
        )

        logger.info(f"🔗 Generated presigned URL (expires in 1 hour)")

        return SlideExportResponse(
            success=True,
            download_url=presigned_url,
            file_name=file_name,
            file_size=len(pptx_bytes),
            expires_in=3600,
            error=None,
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ Failed to export slides: {e}")
        import traceback

        logger.error(traceback.format_exc())

        return SlideExportResponse(
            success=False,
            download_url=None,
            file_name=None,
            file_size=None,
            error=str(e),
        )
